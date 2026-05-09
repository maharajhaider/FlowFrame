import structlog
import hashlib
import os
from typing import List, Dict, Any, Optional, Union, BinaryIO
from io import BytesIO
import aiofiles
import tempfile
import uuid

# Document processing imports
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
from PIL import Image

# LangChain imports for text processing
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document as LangChainDocument

# Local imports
import sys
import os
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from utils.config import Config

logger = structlog.get_logger()

class DocumentProcessor:
    """
    Handles document processing, text extraction, and chunking
    """
    
    def __init__(self):
        self.config = Config()
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.config.CHUNK_SIZE,
            chunk_overlap=self.config.CHUNK_OVERLAP,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        logger.info("DocumentProcessor initialized")
    
    def generate_file_hash(self, content: bytes) -> str:
        """Generate MD5 hash for file content to detect duplicates"""
        return hashlib.md5(content).hexdigest()
    
    def validate_file_type(self, content_type: str) -> bool:
        """Validate if file type is supported"""
        return (content_type in self.config.SUPPORTED_TEXT_TYPES or 
                content_type in self.config.SUPPORTED_IMAGE_TYPES)
    
    def validate_file_size(self, content: bytes, content_type: str) -> bool:
        """Validate file size based on type"""
        size = len(content)
        
        if content_type in self.config.SUPPORTED_IMAGE_TYPES:
            return size <= self.config.get_max_image_size_bytes()
        else:
            return size <= self.config.get_max_file_size_bytes()
    
    async def extract_text_from_pdf(self, content: bytes) -> str:
        """Extract text from PDF file"""
        try:
            # Check for PDF header but don't fail immediately
            has_pdf_header = content.startswith(b'%PDF-')
            if not has_pdf_header:
                logger.warning("Invalid PDF file - missing PDF header, attempting text extraction anyway")
            
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_file:
                temp_file.write(content)
                temp_file.flush()
                temp_file_path = temp_file.name
                
            try:
                text = ""
                with open(temp_file_path, 'rb') as file:
                    try:
                        pdf_reader = PyPDF2.PdfReader(file)
                        
                        # Check if PDF is encrypted
                        if pdf_reader.is_encrypted:
                            logger.warning("PDF is encrypted, cannot extract text")
                            return f"PDF file {temp_file_path} is encrypted and cannot be processed.\nFile size: {len(content)} bytes.\nThis document was uploaded but text extraction failed due to encryption."
                        
                        for page_num, page in enumerate(pdf_reader.pages):
                            try:
                                page_text = page.extract_text()
                                if page_text:
                                    text += page_text + "\n"
                            except Exception as page_error:
                                logger.warning(f"Failed to extract text from page {page_num}", error=str(page_error))
                                continue
                                
                    except Exception as pdf_error:
                        logger.warning(f"Failed to initialize PDF reader: {str(pdf_error)}")
                        # If PDF processing fails completely, return informative fallback
                        return f"PDF processing failed for file.\nError: {str(pdf_error)}\nFile size: {len(content)} bytes.\nThis document was uploaded but could not be processed due to format issues.\nPlease ensure the file is a valid PDF format."
                
                if text.strip():
                    return text.strip()
                else:
                    return f"PDF file appears to be valid but contains no extractable text.\nFile size: {len(content)} bytes.\nThis may be an image-only PDF or contains non-text content.\nDocument has been uploaded to the system."
                
            finally:
                # Always clean up temp file
                try:
                    os.unlink(temp_file_path)
                except Exception:
                    pass  # Ignore cleanup errors
                    
        except Exception as e:
            logger.error("Error extracting text from PDF", error=str(e))
            # Return substantial fallback text instead of empty string
            return f"PDF processing encountered an error: {str(e)}\nFile size: {len(content)} bytes.\nThe document was uploaded but text extraction failed.\nThis content serves as a placeholder to ensure the document is tracked in the system."
    
    async def extract_text_from_docx(self, content: bytes) -> str:
        """Extract text from Word document"""
        try:
            # Check for ZIP header
            has_zip_header = content.startswith(b'PK')
            if not has_zip_header:
                logger.warning("Invalid DOCX file - missing ZIP header")
                # Try to decode as plain text if it's a small file
                if len(content) < 10000:  # Only try for small files
                    try:
                        decoded_text = content.decode('utf-8', errors='ignore').strip()
                        if decoded_text and len(decoded_text) > 10:
                            logger.info("Successfully decoded non-DOCX file as plain text")
                            return f"Plain text content from {content[:100].decode('utf-8', errors='ignore')}...\n\nFull content:\n{decoded_text}"
                    except Exception:
                        pass
                
                # Return informative fallback for invalid DOCX
                return f"Invalid DOCX file format detected.\nFile size: {len(content)} bytes.\nThis file does not have a valid DOCX structure (missing ZIP header).\nPlease ensure the file is properly saved as a .docx format.\nDocument uploaded for reference but content could not be extracted."
            
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as temp_file:
                temp_file.write(content)
                temp_file.flush()
                temp_file_path = temp_file.name
                
            try:
                try:
                    doc = Document(temp_file_path)
                    text_parts = []
                    
                    # Extract text from paragraphs
                    paragraph_count = 0
                    for paragraph in doc.paragraphs:
                        if paragraph.text.strip():
                            text_parts.append(paragraph.text.strip())
                            paragraph_count += 1
                    
                    # Extract text from tables if any
                    table_count = 0
                    for table in doc.tables:
                        table_count += 1
                        for row in table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    text_parts.append(cell.text.strip())
                    
                    logger.info(f"DOCX processing: found {paragraph_count} paragraphs and {table_count} tables")
                    
                    if text_parts:
                        text = "\n".join(text_parts)
                        logger.info(f"Successfully extracted {len(text)} characters from DOCX")
                        return text.strip()
                    else:
                        fallback_text = f"Valid DOCX file but no readable text content found.\nFile size: {len(content)} bytes.\nParagraphs checked: {paragraph_count}\nTables checked: {table_count}\nThis may be an empty document or contains only formatting/images.\nDocument has been uploaded to the system for reference."
                        logger.warning("DOCX file processed but no text content found")
                        return fallback_text
                    
                except Exception as docx_error:
                    logger.warning(f"Failed to initialize Document reader: {str(docx_error)}")
                    fallback_text = f"DOCX processing failed.\nError: {str(docx_error)}\nFile size: {len(content)} bytes.\nThis document was uploaded but could not be processed due to format issues.\nThe file may be corrupted or use an unsupported DOCX variant.\nDocument uploaded for reference."
                    return fallback_text
                
            finally:
                # Always clean up temp file
                try:
                    os.unlink(temp_file_path)
                except Exception:
                    pass  # Ignore cleanup errors
                    
        except Exception as e:
            logger.error("Error extracting text from DOCX", error=str(e))
            # Return substantial fallback text instead of empty string
            fallback_text = f"DOCX processing encountered an unexpected error: {str(e)}\nFile size: {len(content)} bytes.\nThe document was uploaded but text extraction failed.\nThis content serves as a placeholder to ensure the document is tracked in the system and can be referenced later."
            return fallback_text
    
    async def extract_text_from_pptx(self, content: bytes) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            # Check for ZIP header (PPTX is a ZIP archive)
            if not content.startswith(b'PK'):
                logger.warning("Invalid PPTX file - missing ZIP header")
                return ""
            
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                temp_file.write(content)
                temp_file.flush()
                temp_file_path = temp_file.name
                
            try:
                presentation = Presentation(temp_file_path)
                text_runs = []
                
                for slide_num, slide in enumerate(presentation.slides):
                    try:
                        slide_text = f"Slide {slide_num + 1}:\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text += shape.text.strip() + "\n"
                        
                        if slide_text.strip() != f"Slide {slide_num + 1}:":
                            text_runs.append(slide_text)
                    except Exception as slide_error:
                        logger.warning(f"Failed to extract text from slide {slide_num + 1}", error=str(slide_error))
                        continue
                
                return "\n\n".join(text_runs).strip()
                
            finally:
                # Always clean up temp file
                try:
                    os.unlink(temp_file_path)
                except Exception:
                    pass  # Ignore cleanup errors
                    
        except Exception as e:
            logger.error("Error extracting text from PPTX", error=str(e))
            # Return empty string instead of raising to allow processing to continue
            return ""
    
    async def extract_text_from_xlsx(self, content: bytes) -> str:
        """Extract text from Excel file"""
        try:
            # Check for ZIP header (XLSX is a ZIP archive)
            if not content.startswith(b'PK'):
                logger.warning("Invalid XLSX file - missing ZIP header")
                return ""
            
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as temp_file:
                temp_file.write(content)
                temp_file.flush()
                temp_file_path = temp_file.name
                
            try:
                workbook = openpyxl.load_workbook(temp_file_path, read_only=True, data_only=True)
                text_parts = []
                
                for sheet_name in workbook.sheetnames:
                    try:
                        sheet = workbook[sheet_name]
                        sheet_text = f"Sheet: {sheet_name}\n"
                        row_count = 0
                        
                        for row in sheet.iter_rows(values_only=True, max_row=1000):  # Limit rows to prevent memory issues
                            row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                            if row_text.strip():
                                sheet_text += row_text + "\n"
                                row_count += 1
                        
                        if row_count > 0:
                            text_parts.append(sheet_text)
                    except Exception as sheet_error:
                        logger.warning(f"Failed to extract text from sheet {sheet_name}", error=str(sheet_error))
                        continue
                
                workbook.close()
                return "\n\n".join(text_parts).strip()
                
            finally:
                # Always clean up temp file
                try:
                    os.unlink(temp_file_path)
                except Exception:
                    pass  # Ignore cleanup errors
                    
        except Exception as e:
            logger.error("Error extracting text from XLSX", error=str(e))
            # Return empty string instead of raising to allow processing to continue
            return ""
    
    async def process_image(self, content: bytes, filename: str) -> str:
        """Process image file - extract basic metadata"""
        try:
            image = Image.open(BytesIO(content))
            
            # Extract basic image information
            info = {
                "filename": filename,
                "format": image.format,
                "size": f"{image.width}x{image.height}",
                "mode": image.mode
            }
            
            # Try to extract EXIF data if available
            if hasattr(image, '_getexif') and image._getexif():
                exif = image._getexif()
                if exif:
                    info["has_exif"] = True
            
            return f"Image: {filename}\nFormat: {info['format']}\nSize: {info['size']}\nMode: {info['mode']}"
            
        except Exception as e:
            logger.error("Error processing image", error=str(e))
            return f"Image file: {filename} (could not process image details)"
    
    async def extract_text(self, content: bytes, content_type: str, filename: str) -> str:
        """Extract text from various file types"""
        try:
            logger.info(f"Extracting text from {filename} (type: {content_type}, size: {len(content)} bytes)")
            
            # Handle empty files
            if len(content) == 0:
                logger.warning(f"File {filename} is empty")
                return ""
            
            extracted_text = ""
            
            if content_type == 'text/plain' or content_type == 'text/markdown' or filename.endswith('.txt'):
                try:
                    # Try UTF-8 first, then fallback to other encodings
                    try:
                        extracted_text = content.decode('utf-8')
                        logger.info(f"Successfully decoded {filename} as UTF-8")
                    except UnicodeDecodeError:
                        try:
                            extracted_text = content.decode('latin-1')
                            logger.info(f"Successfully decoded {filename} as latin-1")
                        except UnicodeDecodeError:
                            extracted_text = content.decode('utf-8', errors='ignore')
                            logger.warning(f"Decoded {filename} with errors ignored")
                    
                    # Ensure we have meaningful text
                    if extracted_text and len(extracted_text.strip()) > 5:
                        logger.info(f"Text file {filename} contains {len(extracted_text)} characters")
                    else:
                        logger.warning(f"Text file {filename} appears to be empty or very short")
                        
                except Exception as e:
                    logger.warning(f"Failed to decode text file {filename}", error=str(e))
                    extracted_text = f"Failed to decode text file {filename}: {str(e)}\nFile size: {len(content)} bytes\nThis document was uploaded but could not be read as text."
            
            elif content_type == 'application/pdf':
                extracted_text = await self.extract_text_from_pdf(content)
            
            elif content_type in ['application/msword', 
                                'application/vnd.openxmlformats-officedocument.wordprocessingml.document']:
                extracted_text = await self.extract_text_from_docx(content)
            
            elif content_type in ['application/vnd.ms-powerpoint',
                                'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
                extracted_text = await self.extract_text_from_pptx(content)
            
            elif content_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                extracted_text = await self.extract_text_from_xlsx(content)
            
            elif content_type in self.config.SUPPORTED_IMAGE_TYPES:
                extracted_text = await self.process_image(content, filename)
            
            else:
                logger.warning(f"Unsupported content type: {content_type}")
                extracted_text = f"Unsupported file type: {filename} ({content_type})"
            
            # Clean up the extracted text
            if extracted_text:
                extracted_text = extracted_text.strip()
                
            logger.info(f"Extracted {len(extracted_text)} characters from {filename}")
            
            # Return some default content if extraction yielded nothing
            if not extracted_text:
                logger.warning(f"No text extracted from {filename}")
                extracted_text = f"Document: {filename}\nContent type: {content_type}\nFile size: {len(content)} bytes\n(No readable text content found)"
            
            return extracted_text
                
        except Exception as e:
            logger.error("Error extracting text", content_type=content_type, filename=filename, error=str(e))
            # Return a fallback description instead of raising
            return f"Error processing {filename}: {str(e)}"
    
    def chunk_text(self, text: str, metadata: Dict[str, Any] = None) -> List[LangChainDocument]:
        """Split text into chunks using LangChain text splitter"""
        try:
            if not text or not text.strip():
                logger.warning("Empty text provided for chunking")
                return []
            
            text = text.strip()
            text_length = len(text)
            
            logger.info(f"Chunking text with length: {text_length}")
            
            # Always create at least one chunk, even for very short texts
            # This ensures that all processed documents get stored
            if text_length < 500:  # Increased threshold from 200 to 500
                logger.info(f"Text length {text_length} chars - creating single chunk")
                chunk_metadata = {
                    "chunk_index": 0,
                    "total_chunks": 1,
                    "text_length": text_length,
                    "chunk_type": "single_chunk",
                    **(metadata or {})
                }
                
                return [LangChainDocument(
                    page_content=text,
                    metadata=chunk_metadata
                )]
            
            # Split text into chunks
            try:
                chunks = self.text_splitter.split_text(text)
            except Exception as split_error:
                logger.warning(f"Text splitter failed: {str(split_error)}, creating single chunk")
                # Fallback to single chunk if splitter fails
                chunk_metadata = {
                    "chunk_index": 0,
                    "total_chunks": 1,
                    "text_length": text_length,
                    "chunk_type": "fallback_single_chunk",
                    "split_error": str(split_error),
                    **(metadata or {})
                }
                
                return [LangChainDocument(
                    page_content=text,
                    metadata=chunk_metadata
                )]
            
            if not chunks:
                logger.warning("Text splitter returned no chunks, creating fallback chunk")
                # Create a fallback chunk if splitter returns empty
                chunk_metadata = {
                    "chunk_index": 0,
                    "total_chunks": 1,
                    "text_length": text_length,
                    "chunk_type": "fallback_empty_splitter",
                    **(metadata or {})
                }
                
                return [LangChainDocument(
                    page_content=text,
                    metadata=chunk_metadata
                )]
            
            logger.info(f"Text splitter created {len(chunks)} chunks from text")
            
            # Create LangChain documents with metadata
            documents = []
            for i, chunk in enumerate(chunks):
                if chunk.strip():  # Only include non-empty chunks
                    chunk_metadata = {
                        "chunk_index": i,
                        "total_chunks": len(chunks),
                        "chunk_length": len(chunk),
                        "chunk_type": "normal_chunk",
                        **(metadata or {})
                    }
                    
                    documents.append(LangChainDocument(
                        page_content=chunk.strip(),
                        metadata=chunk_metadata
                    ))
            
            # Ensure we always return at least one chunk
            if not documents:
                logger.warning("All chunks were empty after filtering, creating fallback chunk")
                chunk_metadata = {
                    "chunk_index": 0,
                    "total_chunks": 1,
                    "text_length": text_length,
                    "chunk_type": "fallback_empty_filter",
                    **(metadata or {})
                }
                
                documents = [LangChainDocument(
                    page_content=text,
                    metadata=chunk_metadata
                )]
            
            logger.info(f"Created {len(documents)} final document chunks")
            return documents
            
        except Exception as e:
            logger.error("Error chunking text", error=str(e), text_length=len(text) if text else 0)
            # Create an error chunk instead of raising
            error_text = f"Error chunking text: {str(e)}\nOriginal text length: {len(text) if text else 0}\nOriginal text: {text[:500] if text else 'None'}..."
            chunk_metadata = {
                "chunk_index": 0,
                "total_chunks": 1,
                "text_length": len(error_text),
                "chunk_type": "error_chunk",
                "error": str(e),
                **(metadata or {})
            }
            
            return [LangChainDocument(
                page_content=error_text,
                metadata=chunk_metadata
            )]
    
    async def process_document(self, 
                             content: bytes, 
                             filename: str, 
                             content_type: str,
                             collection_type: str = "documents",
                             additional_metadata: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        Complete document processing pipeline
        
        Args:
            content: Raw file content
            filename: Original filename
            content_type: MIME type
            collection_type: Type of collection for categorization
            additional_metadata: Additional metadata to include
            
        Returns:
            Processed document with chunks and metadata
        """
        try:
            # Validate file type (but don't fail hard)
            if not self.validate_file_type(content_type):
                logger.warning(f"Unsupported file type: {content_type}")
                extracted_text = f"Unsupported file: {filename} ({content_type})"
            elif not self.validate_file_size(content, content_type):
                max_size = (self.config.MAX_IMAGE_SIZE_MB if content_type in self.config.SUPPORTED_IMAGE_TYPES 
                           else self.config.MAX_FILE_SIZE_MB)
                logger.warning(f"File too large: {filename} (size: {len(content)} bytes, max: {max_size}MB)")
                extracted_text = f"File too large: {filename} (size: {len(content)} bytes)"
            else:
                # Extract text (this now returns fallback text on errors instead of raising)
                extracted_text = await self.extract_text(content, content_type, filename)
            
            # Generate file hash for duplicate detection
            file_hash = self.generate_file_hash(content)
            
            # Prepare metadata
            base_metadata = {
                "filename": filename,
                "content_type": content_type,
                "file_hash": file_hash,
                "file_size": len(content),
                "collection_type": collection_type,
                "document_id": str(uuid.uuid4()),
                "processed_at": str(os.times().elapsed)
            }
            
            if additional_metadata:
                base_metadata.update(additional_metadata)
            
            # Add more metadata about extraction
            base_metadata["extracted_text_length"] = len(extracted_text)
            base_metadata["total_chunks"] = 0  # Will be updated after chunking
            
            # Chunk the text (this handles empty text gracefully)
            logger.info(f"About to chunk text: length={len(extracted_text)}, first 100 chars: {extracted_text[:100]}")
            chunks = self.chunk_text(extracted_text, base_metadata)
            
            # Update metadata with actual chunk count
            base_metadata["total_chunks"] = len(chunks)
            
            result = {
                "document_id": base_metadata["document_id"],
                "filename": filename,
                "content_type": content_type,
                "file_hash": file_hash,
                "extracted_text": extracted_text,
                "chunks": chunks,
                "metadata": base_metadata,
                "chunk_count": len(chunks)
            }
            
            logger.info(
                "Document processed successfully",
                filename=filename,
                content_type=content_type,
                chunks=len(chunks),
                text_length=len(extracted_text),
                chunk_details=[{
                    "index": i,
                    "length": len(chunk.page_content),
                    "type": chunk.metadata.get("chunk_type", "unknown")
                } for i, chunk in enumerate(chunks[:3])]  # Log first 3 chunks
            )
            
            return result
            
        except Exception as e:
            logger.error("Error processing document", filename=filename, error=str(e))
            
            # Create a fallback result instead of raising
            fallback_text = f"Failed to process {filename}: {str(e)}"
            file_hash = self.generate_file_hash(content) if content else "error"
            
            base_metadata = {
                "filename": filename,
                "content_type": content_type,
                "file_hash": file_hash,
                "file_size": len(content) if content else 0,
                "collection_type": collection_type,
                "document_id": str(uuid.uuid4()),
                "processed_at": str(os.times().elapsed),
                "processing_error": str(e)
            }
            
            if additional_metadata:
                base_metadata.update(additional_metadata)
            
            # Create a single chunk with error information
            chunks = self.chunk_text(fallback_text, base_metadata)
            
            return {
                "document_id": base_metadata["document_id"],
                "filename": filename,
                "content_type": content_type,
                "file_hash": file_hash,
                "extracted_text": fallback_text,
                "chunks": chunks,
                "metadata": base_metadata,
                "chunk_count": len(chunks),
                "processing_error": str(e)
            } 